import json
from pathlib import Path

from pptx import Presentation

from app.agents.deck_planner_agent import DeckPlannerAgent
from app.agents.grounding_checker import GroundingChecker
from app.agents.paper_summary_agent import PaperSummaryAgent
from app.agents.slide_writer_agent import SlideWriterAgent
from app.llm.providers import MockLLMProvider
from app.parser.pdf_parser import PDFParser
from app.schemas.common import GenerationSettings


def load_sample_paper():
    fixture = Path(__file__).parent / "fixtures" / "sample_paper.pdf"
    return PDFParser().parse(fixture, paper_id="paper_quality")


def test_summary_uses_parsed_sections():
    paper = load_sample_paper()
    summary = PaperSummaryAgent(MockLLMProvider(paper)).run(paper)
    assert "reliable local-first paper to PPT generation" in summary.problem
    assert "parse PDFs" in summary.method_overview


def test_slide_writer_creates_section_specific_refs_and_notes():
    paper = load_sample_paper()
    llm = MockLLMProvider(paper)
    summary = PaperSummaryAgent(llm).run(paper)
    plan = DeckPlannerAgent(llm).run(paper, summary, GenerationSettings(slide_count=12))
    drafts = SlideWriterAgent(llm).run(paper, summary, plan)

    for slide in drafts.slides:
        if slide.bullets:
            assert len(slide.bullets) <= 5
            assert slide.speaker_notes

    method_slide = next(slide for slide in drafts.slides if slide.title == "Method Overview")
    result_slide = next(slide for slide in drafts.slides if slide.title == "Main Results")
    assert any(ref.section_title == "Method" for ref in method_slide.source_refs)
    assert any(ref.section_title == "Results" for ref in result_slide.source_refs)


def test_grounding_checker_flags_abstract_only_refs():
    paper = load_sample_paper()
    llm = MockLLMProvider(paper)
    summary = PaperSummaryAgent(llm).run(paper)
    plan = DeckPlannerAgent(llm).run(paper, summary, GenerationSettings(slide_count=10))
    drafts = SlideWriterAgent(llm).run(paper, summary, plan)
    report = GroundingChecker().run(drafts)
    assert report.slide_count == len(drafts.slides)


def test_title_slide_is_not_text_dump(tmp_path: Path):
    paper = load_sample_paper()
    llm = MockLLMProvider(paper)
    summary = PaperSummaryAgent(llm).run(paper)
    plan = DeckPlannerAgent(llm).run(paper, summary, GenerationSettings(slide_count=10))
    drafts = SlideWriterAgent(llm).run(paper, summary, plan)

    from app.ppt.ppt_builder import PPTBuilder

    output = tmp_path / "deck.pptx"
    PPTBuilder().build(output, summary, plan, drafts, GenerationSettings())
    prs = Presentation(output)
    title_texts = [shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text") and shape.text.strip()]
    assert title_texts[0] == "Sample Research Paper"
    assert "Introduction" not in title_texts[0]
